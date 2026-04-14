#!/usr/bin/env python3
"""
Extract plain text from PDF and PPTX under assets/UC/ and write docs/from_intermedia_uc.md.
Requires: .venv-extract with pypdf, python-pptx (see repo README or run:
  python3 -m venv .venv-extract && . .venv-extract/bin/activate && pip install pypdf python-pptx
"""

from __future__ import annotations

try:
    import pypdf  # noqa: F401
    import pptx  # noqa: F401
except ImportError as e:
    raise SystemExit(
        "Missing Python deps (pypdf, python-pptx). From repo root:\n"
        "  python3 -m venv .venv-extract && .venv-extract/bin/pip install pypdf python-pptx\n"
        "  .venv-extract/bin/python scripts/extract_intermedia_uc.py\n"
        f"Original error: {e}"
    ) from e

import re
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

ROOT = Path(__file__).resolve().parents[1]
UC = ROOT / "assets" / "UC"
OUT = ROOT / "docs" / "from_intermedia_uc.md"


def extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            t = page.extract_text() or ""
        except Exception as e:  # noqa: BLE001
            t = f"[extract error page {i + 1}: {e}]"
        parts.append(t)
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        parts: list[str] = []
        for si, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [
                            (cell.text or "").strip() for cell in row.cells
                        ]
                        lines.append(" | ".join(cells))
            if lines:
                parts.append(f"--- Slide {si} ---\n" + "\n".join(lines))
        return "\n\n".join(parts)
    except Exception:
        # Fallback: raw XML text from slide parts (no python-pptx table nuance)
        return extract_pptx_zip_fallback(path)


def extract_pptx_zip_fallback(path: Path) -> str:
    """Minimal text from pptx zip if Presentation fails."""
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    chunks: list[str] = []
    with zipfile.ZipFile(path) as z:
        names = sorted(
            n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )
        for name in names:
            data = z.read(name)
            root = ET.fromstring(data)
            texts = []
            for t in root.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                if t.text:
                    texts.append(t.text)
            if texts:
                slide_id = re.search(r"slide(\d+)", name)
                sid = slide_id.group(1) if slide_id else name
                chunks.append(f"--- Slide {sid} ---\n" + "\n".join(texts))
    return "\n\n".join(chunks) if chunks else "[no text extracted from pptx]"


def normalize_whitespace(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()


def main() -> int:
    if not UC.is_dir():
        print(f"Missing {UC}", file=sys.stderr)
        return 1

    pdfs = sorted(UC.rglob("*.pdf"))
    pptxs = sorted(UC.rglob("*.pptx"))

    lines: list[str] = [
        "# Intermedia UC — extracted text (partner collateral)",
        "",
        "_Auto-generated from PDF and PPTX under `assets/UC/` via `scripts/extract_intermedia_uc.py` "
        "(`npm run docs:extract-intermedia-uc`). See [`from_intermedia_uc_themes.md`](from_intermedia_uc_themes.md) "
        "for folder-level theme summaries (sub-agent review; not overwritten on re-extract). "
        "For internal training reference; verify against current partner agreements and GPC-approved messaging before external use._",
        "",
        f"- **PDF files:** {len(pdfs)}",
        f"- **PPTX files:** {len(pptxs)}",
        "",
        "---",
        "",
    ]

    all_files = [(p, "pdf") for p in pdfs] + [(p, "pptx") for p in pptxs]
    all_files.sort(key=lambda x: str(x[0]).lower())

    if not all_files:
        if OUT.exists():
            print(
                f"No PDF/PPTX under {UC}; leaving existing {OUT} unchanged.",
                file=sys.stderr,
            )
            return 0
        print(
            f"No PDF/PPTX found under {UC}. Add source files or remove this check.",
            file=sys.stderr,
        )
        return 1

    for path, kind in all_files:
        rel = path.relative_to(ROOT)
        lines.append(f"## `{rel}`")
        lines.append("")
        try:
            if kind == "pdf":
                body = extract_pdf(path)
            else:
                body = extract_pptx(path)
            body = normalize_whitespace(body)
            if not body:
                body = "_[No text extracted — may be image-only or encrypted]_"
        except Exception as e:  # noqa: BLE001
            body = f"_Extraction failed: {e}_"
        lines.append(body)
        lines.append("")
        lines.append("---")
        lines.append("")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {OUT} ({OUT.stat().st_size} bytes)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
